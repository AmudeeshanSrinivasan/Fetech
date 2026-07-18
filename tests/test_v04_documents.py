from __future__ import annotations

import asyncio
import io
import json
import zipfile
from collections.abc import Callable
from datetime import datetime
from hashlib import sha256
from pathlib import Path
from typing import Any
from uuid import uuid4

import pytest

from fetech.adapters.base import (
    AdapterBudgetExceededError,
    AdapterDependencyError,
    AdapterExecutionError,
    ExecutionContext,
)
from fetech.adapters.documents import (
    DOCUMENT_CAPABILITIES,
    DocumentAdapter,
    DocumentLimits,
    DocumentParseWorker,
    GitLFSResolvedObject,
    GitLFSResolveRequest,
    PDFOCRPage,
    _detect_capability,
    _parse,
)
from fetech.logic.base import BackendExecutionError
from fetech.logic.process import ProcessResult
from fetech.models import (
    AttemptStatus,
    FetchAttempt,
    FetchRequest,
    PageState,
    PlanNode,
    QualityAssessment,
    Resource,
    ResourceBudget,
)
from fetech.registry import CapabilityRegistry
from fetech.storage import FileSystemCAS, build_artifact


def _pdf_bytes() -> bytes:
    from pypdf import PdfWriter

    output = io.BytesIO()
    writer = PdfWriter()
    writer.add_blank_page(width=72, height=72)
    writer.write(output)
    return output.getvalue()


def _docx_bytes() -> bytes:
    from docx import Document

    output = io.BytesIO()
    document = Document()
    document.add_paragraph("Bounded DOCX fixture content with a stable paragraph locator.")
    document.save(output)
    return output.getvalue()


def _pptx_bytes() -> bytes:
    from pptx import Presentation

    output = io.BytesIO()
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[1])
    slide.shapes.title.text = "Bounded presentation fixture"
    presentation.save(output)
    return output.getvalue()


def _xlsx_bytes() -> bytes:
    from openpyxl import Workbook

    output = io.BytesIO()
    workbook = Workbook()
    sheet = workbook.active
    sheet.append(["name", "created"])
    sheet.append(["fixture", datetime(2026, 7, 17, 1, 2, 3)])
    workbook.save(output)
    workbook.close()
    return output.getvalue()


def _zip_bytes(name: str = "folder/file.txt", content: bytes = b"bounded fixture") -> bytes:
    output = io.BytesIO()
    with zipfile.ZipFile(output, "w", compression=zipfile.ZIP_DEFLATED) as archive:
        archive.writestr(name, content)
    return output.getvalue()


def _lfs_pointer(body: bytes) -> bytes:
    return (
        b"version https://git-lfs.github.com/spec/v1\n"
        b"oid sha256:"
        + sha256(body).hexdigest().encode("ascii")
        + b"\nsize "
        + str(len(body)).encode("ascii")
        + b"\n"
    )


def _identity(value: bytes) -> bytes:
    return value


DOCUMENT_CASES: tuple[
    tuple[str, str, Callable[[], bytes] | bytes, str, str],
    ...,
] = (
    ("pdf", "/report.pdf", _pdf_bytes, "pypdf", "page:1"),
    ("scanned_pdf", "/scan.pdf", _pdf_bytes, "pypdf", "page:1"),
    ("docx", "/report.docx", _docx_bytes, "python-docx", "paragraph:1"),
    ("pptx", "/slides.pptx", _pptx_bytes, "python-pptx", "slide:1"),
    ("xlsx", "/table.xlsx", _xlsx_bytes, "openpyxl", "sheet:Sheet/row:1"),
    ("csv", "/table.csv", b"name,value\nfixture,1\n", "csv", "row:1"),
    ("txt", "/note.txt", b"first line\nsecond line", "text", "line:1"),
    ("markdown", "/note.md", b"# Heading\n\nUseful body", "markdown-text", "line:1"),
    ("json_file", "/data.json", b'{"message":"bounded"}', "stdlib-json", "json:/message"),
    ("xml_file", "/data.xml", b"<root><item>bounded</item></root>", "stdlib-xml", "xpath:/root[1]"),
    ("zip_archive", "/bundle.zip", _zip_bytes, "stdlib-zip", "member:folder/file.txt"),
    (
        "github_raw",
        "/owner/repository/main/README.md",
        b"# Raw GitHub document\n\nBounded fixture",
        "markdown-text",
        "line:1",
    ),
    (
        "git_lfs",
        "/objects/asset.txt",
        b"Resolved Git LFS object bytes from the authorized HTTP acquisition boundary.",
        "text",
        "line:1",
    ),
    (
        "dataset_file",
        "/download.parquet",
        b"PAR1\x00\x01\x02\x03bounded-binary-dataset",
        "binary-metadata",
        "byte:0-30",
    ),
)


@pytest.mark.parametrize(
    ("capability", "target", "body_factory", "expected_parser", "expected_locator"),
    DOCUMENT_CASES,
    ids=[case[0] for case in DOCUMENT_CASES],
)
def test_every_document_capability_has_a_bounded_normalization_path(
    capability: str,
    target: str,
    body_factory: Callable[[], bytes] | bytes,
    expected_parser: str,
    expected_locator: str,
) -> None:
    body = body_factory() if callable(body_factory) else _identity(body_factory)
    document, locators, parser = _parse(capability, body, target=target)

    assert document["blocks"]
    assert expected_locator in locators
    assert parser == expected_parser
    json.dumps(document, allow_nan=False)


@pytest.mark.asyncio
@pytest.mark.parametrize(
    ("capability", "target", "body_factory", "expected_parser", "expected_locator"),
    DOCUMENT_CASES,
    ids=[case[0] for case in DOCUMENT_CASES],
)
async def test_real_isolated_worker_normalizes_every_document_capability(
    capability: str,
    target: str,
    body_factory: Callable[[], bytes] | bytes,
    expected_parser: str,
    expected_locator: str,
) -> None:
    body = body_factory() if callable(body_factory) else _identity(body_factory)
    result = await DocumentParseWorker().parse(
        capability,
        body,
        target=target,
        limits=DocumentLimits(),
        timeout_seconds=5,
    )

    assert result.parser == expected_parser
    assert expected_locator in result.locators


def test_document_matrix_is_the_exact_fourteen_entry_manifest_category() -> None:
    registry = CapabilityRegistry()
    manifest_ids = {entry.id for entry in registry.for_category("documents")}

    assert len(DOCUMENT_CAPABILITIES) == 14
    assert manifest_ids == DOCUMENT_CAPABILITIES
    assert {case[0] for case in DOCUMENT_CASES} == manifest_ids


def test_dataset_and_connectors_prefer_bytes_then_safe_suffix_hints() -> None:
    assert _detect_capability("dataset_file", "/misleading.txt", b"%PDF-1.7") == "pdf"
    assert _detect_capability("github_raw", "/README.md", b"# Heading") == "markdown"
    assert _detect_capability("git_lfs", "/asset.json", b'{"ok":true}') == "json_file"


def test_git_lfs_pointer_requires_a_configured_origin_scoped_connector() -> None:
    pointer = (
        b"version https://git-lfs.github.com/spec/v1\n"
        b"ext-0 example-extension-value\n"
        b"oid sha256:" + b"a" * 64 + b"\n"
        b"size 123\n"
    )

    with pytest.raises(ImportError, match="origin-scoped connector"):
        _parse("git_lfs", pointer, target="/asset.bin")


def test_json_is_strict_depth_bounded_and_has_escaped_pointer_locators() -> None:
    document, locators, _ = _parse("json_file", b'{"a/b":{"~key":1}}')
    assert locators == ("json:/a~1b/~0key",)
    assert document["blocks"][0]["value"] == 1

    with pytest.raises(ValueError, match="duplicate JSON"):
        _parse("json_file", b'{"same":1,"same":2}')
    with pytest.raises(ValueError, match="depth"):
        _parse(
            "json_file",
            b'{"a":{"b":{"c":1}}}',
            limits=DocumentLimits(maximum_depth=1),
        )


def test_xml_rejects_dtd_and_is_depth_bounded_with_stable_xpath_locators() -> None:
    with pytest.raises(ValueError, match="DTD"):
        _parse(
            "xml_file",
            b'<!DOCTYPE root [<!ENTITY x "secret">]><root>&x;</root>',
        )
    with pytest.raises(ValueError, match="depth"):
        _parse(
            "xml_file",
            b"<root><one><two/></one></root>",
            limits=DocumentLimits(maximum_depth=1),
        )

    _, locators, _ = _parse(
        "xml_file",
        b"<root><item/><item/></root>",
    )
    assert locators == (
        "xpath:/root[1]",
        "xpath:/root[1]/item[1]",
        "xpath:/root[1]/item[2]",
    )


def test_csv_and_zip_enforce_block_member_expansion_ratio_and_path_bounds() -> None:
    with pytest.raises(ValueError, match="block limit"):
        _parse(
            "csv",
            b"a\nb\n",
            limits=DocumentLimits(maximum_blocks=1),
        )
    with pytest.raises(ValueError, match="path traversal"):
        _parse("zip_archive", _zip_bytes("../escape.txt"))
    with pytest.raises(ValueError, match="compression ratio"):
        _parse(
            "zip_archive",
            _zip_bytes(content=b"A" * 10_000),
            limits=DocumentLimits(maximum_archive_ratio=1),
        )
    two_members = io.BytesIO()
    with zipfile.ZipFile(two_members, "w") as archive:
        archive.writestr("one.txt", "one")
        archive.writestr("two.txt", "two")
    with pytest.raises(ValueError, match="member limit"):
        _parse(
            "zip_archive",
            two_members.getvalue(),
            limits=DocumentLimits(maximum_archive_members=1),
        )
    duplicate_members = io.BytesIO()
    with zipfile.ZipFile(duplicate_members, "w") as archive:
        archive.writestr("same.txt", "one")
        with pytest.warns(UserWarning, match="Duplicate name"):
            archive.writestr("same.txt", "two")
    with pytest.raises(ValueError, match="duplicate archive"):
        _parse("zip_archive", duplicate_members.getvalue())


def test_office_parsers_reject_forged_zip_containers_before_optional_libraries() -> None:
    forged = _zip_bytes("not-an-office-document.txt")
    for capability in ("docx", "pptx", "xlsx"):
        with pytest.raises(ValueError, match="not a valid"):
            _parse(capability, forged)


@pytest.mark.asyncio
async def test_document_worker_sends_only_path_hint_and_validates_output(
    monkeypatch: pytest.MonkeyPatch,
) -> None:
    captured: dict[str, object] = {}

    async def worker(
        arguments: tuple[str, ...],
        stdin: bytes,
        *,
        timeout_seconds: float,
        memory_mb: int,
        maximum_output_bytes: int,
        maximum_file_bytes: int,
        isolation: object,
    ) -> ProcessResult:
        captured.update(
            {
                "arguments": arguments,
                "stdin": stdin,
                "timeout_seconds": timeout_seconds,
                "memory_mb": memory_mb,
                "maximum_output_bytes": maximum_output_bytes,
                "maximum_file_bytes": maximum_file_bytes,
                "isolation": isolation,
            }
        )
        return ProcessResult(
            returncode=0,
            stdout=json.dumps(
                {
                    "document": {
                        "type": "text",
                        "blocks": [{"locator": "line:1", "text": "safe"}],
                    },
                    "locators": ["line:1"],
                    "parser": "text",
                    "observed_capability": "txt",
                }
            ).encode(),
            stderr=b"",
        )

    monkeypatch.setattr("fetech.adapters.documents.run_bounded", worker)
    result = await DocumentParseWorker().parse(
        "txt",
        b"safe",
        target="https://user:password@example.com/private/file.txt?token=never-store",
        limits=DocumentLimits(),
        timeout_seconds=1,
    )
    payload = json.loads(captured["stdin"])

    assert result.parser == "text"
    assert captured["maximum_file_bytes"] == 50_065_536
    assert payload["target_path"] == "file.txt"
    assert b"private" not in captured["stdin"]
    assert b"user" not in captured["stdin"]
    assert b"password" not in captured["stdin"]
    assert b"never-store" not in captured["stdin"]
    assert captured["arguments"][-1] == "fetech.document_worker"
    assert captured["isolation"].profile.value == "document_parser"  # type: ignore[union-attr]


@pytest.mark.asyncio
async def test_document_worker_timeout_and_malformed_output_are_typed_and_sanitized(
    monkeypatch: pytest.MonkeyPatch,
) -> None:
    async def timeout(*args: object, **kwargs: object) -> ProcessResult:
        del args, kwargs
        raise BackendExecutionError("private worker timeout detail")

    monkeypatch.setattr("fetech.adapters.documents.run_bounded", timeout)
    with pytest.raises(AdapterExecutionError, match="bounded document parse process failed") as caught:
        await DocumentParseWorker().parse(
            "txt",
            b"safe",
            target="/file.txt",
            limits=DocumentLimits(),
            timeout_seconds=0.1,
        )
    assert "private worker timeout detail" not in str(caught.value)

    async def malformed(*args: object, **kwargs: object) -> ProcessResult:
        del args, kwargs
        return ProcessResult(
            returncode=0,
            stdout=b"not-json",
            stderr=b"private parser detail",
        )

    monkeypatch.setattr("fetech.adapters.documents.run_bounded", malformed)
    with pytest.raises(AdapterExecutionError, match="malformed output") as caught:
        await DocumentParseWorker().parse(
            "txt",
            b"safe",
            target="/file.txt",
            limits=DocumentLimits(),
            timeout_seconds=1,
        )
    assert "private parser detail" not in str(caught.value)


@pytest.mark.asyncio
async def test_document_worker_cannot_override_deterministic_format_or_locator_contract(
    monkeypatch: pytest.MonkeyPatch,
) -> None:
    async def contradictory(*args: object, **kwargs: object) -> ProcessResult:
        del args, kwargs
        return ProcessResult(
            returncode=0,
            stdout=json.dumps(
                {
                    "document": {
                        "type": "json",
                        "blocks": [{"locator": "json:/forged", "value": True}],
                    },
                    "locators": ["json:/forged"],
                    "parser": "stdlib-json",
                    "observed_capability": "json_file",
                }
            ).encode(),
            stderr=b"",
        )

    monkeypatch.setattr("fetech.adapters.documents.run_bounded", contradictory)
    with pytest.raises(AdapterExecutionError, match="contradicts deterministic routing"):
        await DocumentParseWorker().parse(
            "txt",
            b"plain text",
            target="/file.txt",
            limits=DocumentLimits(),
            timeout_seconds=1,
        )


@pytest.mark.asyncio
async def test_document_worker_dependency_failure_is_typed_without_stderr_leak(
    monkeypatch: pytest.MonkeyPatch,
) -> None:
    async def missing(*args: object, **kwargs: object) -> ProcessResult:
        del args, kwargs
        return ProcessResult(
            returncode=2,
            stdout=b'{"error":"dependency_missing"}',
            stderr=b"private connector detail",
        )

    monkeypatch.setattr("fetech.adapters.documents.run_bounded", missing)
    with pytest.raises(AdapterDependencyError, match=r"fetech\[documents\]") as caught:
        await DocumentParseWorker().parse(
            "pdf",
            _pdf_bytes(),
            target="/scan.pdf",
            limits=DocumentLimits(),
            timeout_seconds=1,
        )
    assert "private connector detail" not in str(caught.value)


async def _acquired_context(
    tmp_path: Path,
    *,
    target: str,
    body: bytes,
    budget: ResourceBudget | None = None,
    canonical_url: str | None = None,
) -> ExecutionContext:
    cas = FileSystemCAS(tmp_path / "cas")
    resource = Resource(
        canonical_url=canonical_url or target,
        requested_url=target,
        authority_url=target,
        media_type="application/pdf",
        status_code=200,
    )
    uri, digest, size = await cas.put(body)
    raw = build_artifact(
        role="source",
        representation="raw",
        media_type="application/pdf",
        cas_uri=uri,
        digest=digest,
        size=size,
        resource=resource,
        extractor="httpx/test",
        quality=QualityAssessment(accepted=True, score=1, completeness=1),
    )
    return ExecutionContext(
        run_id=uuid4(),
        request=FetchRequest(
            target=target,
            budget=budget or ResourceBudget(deadline_seconds=10),
        ),
        cas=cas,
        resources=[resource],
        artifacts=[raw],
    )


@pytest.mark.asyncio
async def test_production_adapter_uses_real_worker_and_preserves_needs_ocr(
    tmp_path: Path,
) -> None:
    context = await _acquired_context(
        tmp_path,
        target="https://publisher.example/scanned.pdf",
        body=_pdf_bytes(),
    )

    await DocumentAdapter().execute(
        PlanNode(
            id="document",
            capability_id="scanned_pdf",
            adapter="documents",
        ),
        context,
    )

    artifact = context.artifacts[-1]
    assert artifact.representation == "document"
    assert artifact.role == "checked-only"
    assert artifact.quality.page_state == PageState.NEEDS_OCR
    assert not artifact.quality.accepted
    assert artifact.locators == ("page:1",)
    assert context.attempts[-1].status == AttemptStatus.SUCCEEDED
    assert context.attempts[-1].parser == "pypdf"
    assert not context.accepted


@pytest.mark.asyncio
async def test_adapter_rejects_body_above_worker_budget_before_subprocess(
    tmp_path: Path,
    monkeypatch: pytest.MonkeyPatch,
) -> None:
    context = await _acquired_context(
        tmp_path,
        target="https://publisher.example/file.txt",
        body=b"too large",
        budget=ResourceBudget(bytes=4, decompressed_bytes=4),
    )
    called = False

    async def should_not_run(*args: object, **kwargs: object) -> ProcessResult:
        nonlocal called
        del args, kwargs
        called = True
        raise AssertionError("worker must not run")

    monkeypatch.setattr("fetech.adapters.documents.run_bounded", should_not_run)
    with pytest.raises(AdapterExecutionError, match="parser byte limit"):
        await DocumentAdapter().execute(
            PlanNode(id="document", capability_id="txt", adapter="documents"),
            context,
        )
    assert not called
    assert context.attempts[-1].failure_code == "document_error"


@pytest.mark.asyncio
async def test_git_lfs_resolver_is_exact_origin_sanitized_bounded_and_lineaged(
    tmp_path: Path,
) -> None:
    resolved_body = (
        b"Useful Git LFS object content with enough deterministic text for "
        b"acceptance, stable provenance, and bounded parsing. "
    ) * 3
    requests: list[GitLFSResolveRequest] = []

    class Resolver:
        async def resolve(
            self,
            request: GitLFSResolveRequest,
        ) -> GitLFSResolvedObject:
            requests.append(request)
            return GitLFSResolvedObject(
                origin="https://git.example",
                body=resolved_body,
            )

    context = await _acquired_context(
        tmp_path,
        target="https://git.example/owner/repo/asset.bin?token=private-value",
        body=_lfs_pointer(resolved_body),
    )
    raw = context.artifacts[-1]

    await DocumentAdapter(git_lfs_resolver=Resolver()).execute(
        PlanNode(id="document", capability_id="git_lfs", adapter="documents"),
        context,
    )

    assert len(requests) == 1
    request = requests[0]
    assert request.origin == "https://git.example"
    assert "private-value" not in request.sanitized_target
    assert "REDACTED" in request.sanitized_target
    assert request.pointer.oid_sha256 == sha256(resolved_body).hexdigest()
    assert request.pointer.size == len(resolved_body)
    assert request.maximum_bytes <= context.request.budget.bytes
    artifact = context.artifacts[-1]
    assert artifact.role == "primary"
    assert raw.artifact_id in artifact.parent_artifact_ids
    assert artifact.quality.accepted
    assert context.attempts[-1].bytes_received == len(resolved_body)
    assert context.attempts[-1].consumed_budget == {
        "bytes": len(resolved_body),
        "decompressed_bytes": len(resolved_body) + artifact.size,
    }
    outcome = next(
        item for item in context.capability_outcomes if item.capability_id == "git_lfs"
    )
    assert outcome.details["git_lfs_resolved"] is True


@pytest.mark.asyncio
async def test_git_lfs_pointer_without_resolver_remains_dependency_missing(
    tmp_path: Path,
) -> None:
    body = b"bounded Git LFS object"
    context = await _acquired_context(
        tmp_path,
        target="https://git.example/owner/repo/asset.bin",
        body=_lfs_pointer(body),
    )

    with pytest.raises(AdapterDependencyError, match="exact-origin resolver"):
        await DocumentAdapter().execute(
            PlanNode(id="document", capability_id="git_lfs", adapter="documents"),
            context,
        )

    assert context.attempts[-1].failure_code == "dependency_missing"


@pytest.mark.asyncio
async def test_git_lfs_malformed_pointer_fails_before_resolver(
    tmp_path: Path,
) -> None:
    called = False

    class Resolver:
        async def resolve(
            self,
            request: GitLFSResolveRequest,
        ) -> GitLFSResolvedObject:
            nonlocal called
            del request
            called = True
            raise AssertionError("malformed pointer must not reach resolver")

    context = await _acquired_context(
        tmp_path,
        target="https://git.example/owner/repo/asset.bin",
        body=(
            b"version https://git-lfs.github.com/spec/v1\n"
            b"oid sha256:not-a-digest\n"
            b"size 42\n"
        ),
    )

    with pytest.raises(AdapterExecutionError, match="pointer grammar"):
        await DocumentAdapter(git_lfs_resolver=Resolver()).execute(
            PlanNode(id="document", capability_id="git_lfs", adapter="documents"),
            context,
        )
    assert not called


@pytest.mark.parametrize(
    ("resolved_origin", "resolved_body", "message"),
    [
        ("https://objects.example", b"expected", "exact-origin"),
        ("https://git.example", b"changed!", "hash mismatch"),
        ("https://git.example", b"expected-extra", "size mismatch"),
    ],
)
@pytest.mark.asyncio
async def test_git_lfs_rejects_wrong_origin_hash_and_size(
    tmp_path: Path,
    resolved_origin: str,
    resolved_body: bytes,
    message: str,
) -> None:
    expected = b"expected"

    class Resolver:
        async def resolve(
            self,
            request: GitLFSResolveRequest,
        ) -> GitLFSResolvedObject:
            del request
            return GitLFSResolvedObject(
                origin=resolved_origin,
                body=resolved_body,
            )

    context = await _acquired_context(
        tmp_path,
        target="https://git.example/owner/repo/asset.bin",
        body=_lfs_pointer(expected),
    )

    with pytest.raises(AdapterExecutionError, match=message):
        await DocumentAdapter(git_lfs_resolver=Resolver()).execute(
            PlanNode(id="document", capability_id="git_lfs", adapter="documents"),
            context,
        )


@pytest.mark.asyncio
async def test_invalid_git_lfs_body_stays_charged_and_is_not_reacquired(
    tmp_path: Path,
) -> None:
    expected = b"expected"
    acquired = b"changed!"
    pointer = _lfs_pointer(expected)
    calls = 0

    class Resolver:
        async def resolve(
            self,
            request: GitLFSResolveRequest,
        ) -> GitLFSResolvedObject:
            nonlocal calls
            calls += 1
            return GitLFSResolvedObject(origin=request.origin, body=acquired)

    budget = ResourceBudget(
        bytes=len(pointer) + len(acquired),
        decompressed_bytes=len(pointer) + len(acquired),
    )
    context = await _acquired_context(
        tmp_path,
        target="https://git.example/owner/repo/asset.bin",
        body=pointer,
        budget=budget,
    )
    context.attempts.append(
        FetchAttempt(
            capability_id="http_get",
            sanitized_destination=context.request.target,
            status=AttemptStatus.SUCCEEDED,
            consumed_budget={
                "bytes": len(pointer),
                "decompressed_bytes": len(pointer),
            },
        )
    )
    adapter = DocumentAdapter(git_lfs_resolver=Resolver())
    node = PlanNode(id="document", capability_id="git_lfs", adapter="documents")

    with pytest.raises(AdapterExecutionError, match="hash mismatch"):
        await adapter.execute(node, context)

    assert context.attempts[-1].consumed_budget == {
        "bytes": len(acquired),
        "decompressed_bytes": len(acquired),
    }
    with pytest.raises(AdapterBudgetExceededError, match="decompressed_bytes"):
        await adapter.execute(node, context)
    assert calls == 1


@pytest.mark.asyncio
async def test_git_lfs_declared_size_and_timeout_are_enforced(
    tmp_path: Path,
) -> None:
    called = False
    oversized = b"x" * 512

    class Resolver:
        async def resolve(
            self,
            request: GitLFSResolveRequest,
        ) -> GitLFSResolvedObject:
            nonlocal called
            called = True
            await asyncio.sleep(0.05)
            return GitLFSResolvedObject(origin=request.origin, body=oversized)

    budget_context = await _acquired_context(
        tmp_path / "budget",
        target="https://git.example/owner/repo/asset.bin",
        body=_lfs_pointer(oversized),
        budget=ResourceBudget(bytes=200, decompressed_bytes=200),
    )
    with pytest.raises(AdapterExecutionError, match="declared size"):
        await DocumentAdapter(git_lfs_resolver=Resolver()).execute(
            PlanNode(id="document", capability_id="git_lfs", adapter="documents"),
            budget_context,
        )
    assert not called

    timeout_body = b"useful bounded object"
    timeout_context = await _acquired_context(
        tmp_path / "timeout",
        target="https://git.example/owner/repo/asset.bin",
        body=_lfs_pointer(timeout_body),
    )
    with pytest.raises(AdapterExecutionError, match="timed out"):
        await DocumentAdapter(
            git_lfs_resolver=Resolver(),
            provider_timeout_seconds=0.001,
        ).execute(
            PlanNode(id="document", capability_id="git_lfs", adapter="documents"),
            timeout_context,
        )


@pytest.mark.parametrize("mode", ["wrong_type", "exception"])
@pytest.mark.asyncio
async def test_git_lfs_provider_failures_are_sanitized(
    tmp_path: Path,
    mode: str,
) -> None:
    body = b"bounded object"

    class Resolver:
        async def resolve(self, request: GitLFSResolveRequest) -> Any:
            del request
            if mode == "exception":
                raise RuntimeError("private resolver credential detail")
            return {"origin": "https://git.example", "body": body}

    context = await _acquired_context(
        tmp_path,
        target="https://git.example/owner/repo/asset.bin",
        body=_lfs_pointer(body),
    )

    with pytest.raises(AdapterExecutionError) as caught:
        await DocumentAdapter(git_lfs_resolver=Resolver()).execute(
            PlanNode(id="document", capability_id="git_lfs", adapter="documents"),
            context,
        )
    assert "private resolver credential detail" not in str(caught.value)
    assert context.attempts[-1].failure_code == "document_error"


@pytest.mark.parametrize(
    "target",
    [
        "https://example.com/owner/repo/main/file.txt",
        "https://raw.githubusercontent.com/owner/repo",
        "https://raw.githubusercontent.com/owner/repo/main/%2e%2e/secret.txt",
        "https://raw.githubusercontent.com/owner/repo/main/file.txt?token=value",
    ],
)
@pytest.mark.asyncio
async def test_github_raw_rejects_wrong_origin_and_malformed_paths(
    tmp_path: Path,
    target: str,
) -> None:
    context = await _acquired_context(
        tmp_path,
        target=target,
        body=b"Useful bounded raw GitHub content with deterministic text. " * 4,
    )

    with pytest.raises(AdapterExecutionError, match="GitHub raw target"):
        await DocumentAdapter().execute(
            PlanNode(id="document", capability_id="github_raw", adapter="documents"),
            context,
        )


@pytest.mark.asyncio
async def test_github_raw_accepts_only_the_exact_raw_origin_and_file_path(
    tmp_path: Path,
) -> None:
    context = await _acquired_context(
        tmp_path,
        target="https://raw.githubusercontent.com/owner/repo/main/README.md",
        body=b"Useful bounded raw GitHub content with deterministic text. " * 4,
    )

    await DocumentAdapter().execute(
        PlanNode(id="document", capability_id="github_raw", adapter="documents"),
        context,
    )

    assert context.artifacts[-1].quality.accepted
    assert context.attempts[-1].parser == "markdown-text"


@pytest.mark.asyncio
async def test_github_raw_rejects_a_cross_origin_acquired_resource(
    tmp_path: Path,
) -> None:
    target = "https://raw.githubusercontent.com/owner/repo/main/README.md"
    context = await _acquired_context(
        tmp_path,
        target=target,
        canonical_url="https://example.com/redirected/README.md",
        body=b"Useful bounded redirected content with deterministic text. " * 4,
    )

    with pytest.raises(
        AdapterExecutionError,
        match=r"exact raw\.githubusercontent\.com",
    ):
        await DocumentAdapter().execute(
            PlanNode(id="document", capability_id="github_raw", adapter="documents"),
            context,
        )


@pytest.mark.asyncio
async def test_configured_pdf_ocr_transitions_needs_ocr_to_accepted_with_lineage(
    tmp_path: Path,
) -> None:
    calls: list[dict[str, object]] = []

    class OCR:
        async def extract_pdf(
            self,
            body: bytes,
            *,
            page_count: int,
            language: str | None,
            timeout_seconds: float,
            maximum_output_bytes: int,
        ) -> tuple[PDFOCRPage, ...]:
            calls.append(
                {
                    "body": body,
                    "page_count": page_count,
                    "language": language,
                    "timeout_seconds": timeout_seconds,
                    "maximum_output_bytes": maximum_output_bytes,
                }
            )
            return (
                PDFOCRPage(
                    locator="page:1",
                    text=(
                        "Bounded OCR text recovered from page one with enough useful "
                        "content for acceptance and stable source provenance. "
                    )
                    * 3,
                ),
            )

    pdf = _pdf_bytes()
    context = await _acquired_context(
        tmp_path,
        target="https://publisher.example/scanned.pdf",
        body=pdf,
        budget=ResourceBudget(
            deadline_seconds=10,
            bytes=10_000,
            decompressed_bytes=2_000,
        ),
    )
    raw = context.artifacts[-1]

    await DocumentAdapter(pdf_ocr_provider=OCR()).execute(
        PlanNode(id="document", capability_id="scanned_pdf", adapter="documents"),
        context,
    )

    assert len(calls) == 1
    assert calls[0]["body"] == pdf
    assert calls[0]["page_count"] == 1
    assert calls[0]["maximum_output_bytes"] == 2_000
    artifact = context.artifacts[-1]
    assert artifact.role == "primary"
    assert artifact.quality.page_state == PageState.OK
    assert artifact.quality.accepted
    assert artifact.locators == ("page:1",)
    assert raw.artifact_id in artifact.parent_artifact_ids
    assert context.attempts[-1].parser == "configured-pdf-ocr"
    assert context.attempts[-1].consumed_budget == {
        "decompressed_bytes": artifact.size
    }
    assert context.accepted
    assert context.pending_events[-1][0] == "document_ocr_completed"


@pytest.mark.parametrize(
    "pages",
    [
        {"page:1": "not a typed page"},
        (PDFOCRPage(locator="page:2", text="out of range"),),
        (
            PDFOCRPage(locator="page:1", text="first"),
            PDFOCRPage(locator="page:1", text="duplicate"),
        ),
        (PDFOCRPage(locator="page:1", text=""),),
    ],
)
@pytest.mark.asyncio
async def test_pdf_ocr_rejects_invalid_provider_schema(
    tmp_path: Path,
    pages: object,
) -> None:
    class OCR:
        async def extract_pdf(self, *_: object, **__: object) -> Any:
            return pages

    context = await _acquired_context(
        tmp_path,
        target="https://publisher.example/scanned.pdf",
        body=_pdf_bytes(),
    )

    with pytest.raises(AdapterExecutionError, match="invalid"):
        await DocumentAdapter(pdf_ocr_provider=OCR()).execute(
            PlanNode(id="document", capability_id="scanned_pdf", adapter="documents"),
            context,
        )


@pytest.mark.asyncio
async def test_pdf_ocr_output_budget_and_timeout_are_enforced(
    tmp_path: Path,
) -> None:
    class OversizedOCR:
        async def extract_pdf(self, *_: object, **__: object) -> tuple[PDFOCRPage, ...]:
            return (PDFOCRPage(locator="page:1", text="x" * 5_000),)

    budget_context = await _acquired_context(
        tmp_path / "budget",
        target="https://publisher.example/scanned.pdf",
        body=_pdf_bytes(),
        budget=ResourceBudget(bytes=10_000, decompressed_bytes=1_000),
    )
    with pytest.raises(AdapterExecutionError, match="output byte budget"):
        await DocumentAdapter(pdf_ocr_provider=OversizedOCR()).execute(
            PlanNode(id="document", capability_id="scanned_pdf", adapter="documents"),
            budget_context,
        )

    class SlowOCR:
        async def extract_pdf(self, *_: object, **__: object) -> tuple[PDFOCRPage, ...]:
            await asyncio.sleep(0.05)
            return (PDFOCRPage(locator="page:1", text="eventual text"),)

    timeout_context = await _acquired_context(
        tmp_path / "timeout",
        target="https://publisher.example/scanned.pdf",
        body=_pdf_bytes(),
    )
    with pytest.raises(AdapterExecutionError, match="timed out"):
        await DocumentAdapter(
            pdf_ocr_provider=SlowOCR(),
            provider_timeout_seconds=0.001,
        ).execute(
            PlanNode(id="document", capability_id="scanned_pdf", adapter="documents"),
            timeout_context,
        )


@pytest.mark.asyncio
async def test_pdf_ocr_provider_exception_is_sanitized(
    tmp_path: Path,
) -> None:
    class OCR:
        async def extract_pdf(self, *_: object, **__: object) -> tuple[PDFOCRPage, ...]:
            raise RuntimeError("private OCR connector detail")

    context = await _acquired_context(
        tmp_path,
        target="https://publisher.example/scanned.pdf",
        body=_pdf_bytes(),
    )

    with pytest.raises(AdapterExecutionError) as caught:
        await DocumentAdapter(pdf_ocr_provider=OCR()).execute(
            PlanNode(id="document", capability_id="scanned_pdf", adapter="documents"),
            context,
        )
    assert "private OCR connector detail" not in str(caught.value)
    assert context.attempts[-1].failure_code == "document_error"
