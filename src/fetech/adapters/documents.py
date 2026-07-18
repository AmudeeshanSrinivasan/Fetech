"""Bounded document parsing with stable source locators.

The adapter intentionally performs no network or credential-store access.  It
normalizes an already acquired ``raw`` artifact and keeps optional format
engines behind explicit dependency failures.  Archive and structured parsers
are bounded before their output is admitted to the artifact store.
"""

from __future__ import annotations

import asyncio
import base64
import csv
import hmac
import io
import json
import math
import re
import stat
import sys
import zipfile
from collections.abc import Sequence
from dataclasses import dataclass, field
from datetime import UTC, date, datetime
from hashlib import sha256
from pathlib import Path, PurePosixPath
from typing import Any, Protocol
from urllib.parse import unquote_to_bytes, urlsplit
from xml.etree import ElementTree

from fetech.adapters.base import (
    AdapterBudgetExceededError,
    AdapterDependencyError,
    AdapterExecutionError,
    ExecutionContext,
)
from fetech.auth import canonical_origin
from fetech.docling_artifacts import (
    DOCLING_COMPONENT_VERSIONS,
    DOCLING_SLIM_VERSION,
    docling_artifact_root_identity,
    inspect_docling_artifact_bundle,
    verify_docling_artifact_bundle,
)
from fetech.logic.base import LogicBackendError
from fetech.logic.process import run_bounded
from fetech.models import (
    AttemptStatus,
    CapabilityOutcomeStatus,
    FetchAttempt,
    PageState,
    PlanNode,
    QualityAssessment,
)
from fetech.quality import assess_text
from fetech.security import sanitize_url_for_request
from fetech.storage import build_artifact
from fetech.worker_isolation import (
    WorkerIsolationProfile,
    WorkerIsolationRuntime,
)

DOCUMENT_CAPABILITIES = frozenset(
    {
        "pdf",
        "scanned_pdf",
        "docx",
        "pptx",
        "xlsx",
        "csv",
        "txt",
        "markdown",
        "json_file",
        "xml_file",
        "zip_archive",
        "github_raw",
        "git_lfs",
        "dataset_file",
    }
)

_JSON_XML_PREFIX = re.compile(rb"^\s*(?:\{|\[|<\?xml\b|<[A-Za-z_])", re.IGNORECASE)
_FORBIDDEN_XML = re.compile(rb"<!\s*(?:DOCTYPE|ENTITY)\b", re.IGNORECASE)
_GIT_LFS_POINTER = re.compile(
    rb"\Aversion https://git-lfs\.github\.com/spec/v1\r?\n"
    rb"(?:ext-[A-Za-z0-9-]+ [^\r\n]{1,1024}\r?\n)*"
    rb"oid sha256:(?P<oid>[0-9a-f]{64})\r?\n"
    rb"size (?P<size>0|[1-9][0-9]*)\r?\n?\Z"
)
_GIT_LFS_POINTER_PREFIX = b"version https://git-lfs.github.com/spec/v1"
_MAX_GIT_LFS_POINTER_BYTES = 8_192
_MAX_PROVIDER_TARGET_BYTES = 8_192
_MAX_DOCLING_ARTIFACTS_PATH_BYTES = 8_192
_PDF_PAGE_LOCATOR = re.compile(r"page:(?P<page>[1-9][0-9]*)\Z")
_GITHUB_RAW_ORIGIN = "https://raw.githubusercontent.com"
_GITHUB_PATH_COMPONENT = re.compile(r"[A-Za-z0-9_.-]{1,255}\Z")
_TEXT_SUFFIXES = {".text", ".txt"}
_MARKDOWN_SUFFIXES = {".markdown", ".md", ".mdown", ".mkd"}
_ZIP_SIGNATURES = (b"PK\x03\x04", b"PK\x05\x06", b"PK\x07\x08")
_MAX_DOCUMENT_INPUT_BYTES = 50_000_000
_MAX_DOCUMENT_OUTPUT_BYTES = 50_000_000
_DOCLING_SLIM_VERSION = DOCLING_SLIM_VERSION
_DOCLING_REQUIRED_COMPONENT_VERSIONS = dict(DOCLING_COMPONENT_VERSIONS)
_DOCLING_COMPONENT_KEYS = frozenset(
    {
        "docling",
        "docling-core",
        "docling-ibm-models",
        "docling-parse",
        "docling-slim",
    }
)
_DOCLING_REQUIRED_COMPONENT_KEYS = frozenset(
    _DOCLING_REQUIRED_COMPONENT_VERSIONS
)
_DOCLING_FALLBACK_REASONS = frozenset(
    {
        "docling_contract_invalid",
        "docling_conversion_error",
        "docling_incomplete_pages",
        "docling_non_success",
        "docling_parse_failed",
        "docling_timeout",
        "docling_unavailable",
        "docling_version_mismatch",
    }
)
_SAFE_COMPONENT_VERSION = re.compile(r"[A-Za-z0-9][A-Za-z0-9.+_-]{0,63}\Z")
_SHA256_HEX = re.compile(r"[0-9a-f]{64}\Z")
_DOCUMENT_RESULT_CONTRACTS = {
    "csv": ("csv", "table", "row:"),
    "dataset_file": ("binary-metadata", "binary_dataset", "byte:"),
    "docx": ("python-docx", "docx", "paragraph:"),
    "json_file": ("stdlib-json", "json", "json:"),
    "markdown": ("markdown-text", "markdown", "line:"),
    "pdf": ("pypdf", "pdf", "page:"),
    "pptx": ("python-pptx", "pptx", "slide:"),
    "scanned_pdf": ("pypdf", "pdf", "page:"),
    "txt": ("text", "text", "line:"),
    "xlsx": ("openpyxl", "xlsx", "sheet:"),
    "xml_file": ("stdlib-xml", "xml", "xpath:"),
    "zip_archive": ("stdlib-zip", "zip_archive", "member:"),
}


@dataclass(frozen=True, slots=True)
class DocumentLimits:
    """Parser-local ceilings derived from the Python-owned request budget."""

    maximum_input_bytes: int = 16_000_000
    maximum_output_bytes: int = 50_000_000
    maximum_blocks: int = 10_000
    maximum_depth: int = 64
    maximum_archive_members: int = 1_000
    maximum_archive_ratio: float = 100.0

    def __post_init__(self) -> None:
        numeric = (
            self.maximum_input_bytes,
            self.maximum_output_bytes,
            self.maximum_blocks,
            self.maximum_depth,
            self.maximum_archive_members,
            self.maximum_archive_ratio,
        )
        if min(numeric) <= 0 or not math.isfinite(self.maximum_archive_ratio):
            raise ValueError("document parser limits must be positive")


@dataclass(frozen=True, slots=True)
class DocumentParseResult:
    document: dict[str, Any]
    locators: tuple[str, ...]
    parser: str
    observed_capability: str
    parser_components: tuple[tuple[str, str], ...] = ()
    artifact_bundle_id: str | None = None
    fallback_reason: str | None = None


@dataclass(frozen=True, slots=True)
class GitLFSPointer:
    """Validated Git LFS v1 object identity."""

    oid_sha256: str
    size: int


@dataclass(frozen=True, slots=True)
class GitLFSResolveRequest:
    """Sanitized, bounded input for an exact-origin Git LFS resolver."""

    sanitized_target: str
    origin: str
    pointer: GitLFSPointer
    maximum_bytes: int
    timeout_seconds: float


@dataclass(frozen=True, slots=True)
class GitLFSResolvedObject:
    """Untrusted resolver result revalidated by the Python authority."""

    origin: str
    body: bytes = field(repr=False)


class GitLFSResolver(Protocol):
    """Resolve a pointer without receiving credentials or an unsanitized URL."""

    async def resolve(
        self,
        request: GitLFSResolveRequest,
    ) -> GitLFSResolvedObject: ...


@dataclass(frozen=True, slots=True)
class PDFOCRPage:
    """One page-located OCR result returned by a configured provider."""

    locator: str
    text: str = field(repr=False)


class PDFOCRProvider(Protocol):
    """OCR already acquired PDF bytes under explicit time/output bounds."""

    async def extract_pdf(
        self,
        body: bytes,
        *,
        page_count: int,
        language: str | None,
        timeout_seconds: float,
        maximum_output_bytes: int,
    ) -> Sequence[PDFOCRPage]: ...


class DocumentParser(Protocol):
    async def parse(
        self,
        capability: str,
        body: bytes,
        *,
        target: str,
        limits: DocumentLimits,
        timeout_seconds: float,
    ) -> DocumentParseResult: ...


class DocumentParseWorker:
    """Ephemeral, CPU/time/memory/output-bounded document parser process."""

    def __init__(
        self,
        *,
        memory_mb: int = 512,
        docling_memory_mb: int = 4_096,
        docling_artifacts_path: Path | None = None,
        docling_artifacts_sha256: str | None = None,
        isolation: WorkerIsolationRuntime | None = None,
    ) -> None:
        if (
            not isinstance(memory_mb, int)
            or isinstance(memory_mb, bool)
            or memory_mb <= 0
        ):
            raise ValueError("document worker memory limit must be positive")
        if (
            not isinstance(docling_memory_mb, int)
            or isinstance(docling_memory_mb, bool)
            or not 1_024 <= docling_memory_mb <= 8_192
        ):
            raise ValueError(
                "Docling worker memory limit must be between 1024 and 8192 MiB"
            )
        self.memory_mb = memory_mb
        self.docling_memory_mb = docling_memory_mb
        self.docling_artifacts_path = _validate_docling_artifacts_path(
            docling_artifacts_path
        )
        if (self.docling_artifacts_path is None) != (
            docling_artifacts_sha256 is None
        ):
            raise ValueError(
                "Docling artifacts path and expected SHA-256 must be configured together"
            )
        self.isolation = isolation or WorkerIsolationRuntime.from_environment()
        self.docling_artifact_bundle_id = (
            verify_docling_artifact_bundle(
                self.docling_artifacts_path,
                expected_sha256=docling_artifacts_sha256 or "",
            ).bundle_sha256
            if self.docling_artifacts_path is not None
            else None
        )
        self.docling_artifact_root_id = (
            docling_artifact_root_identity(self.docling_artifacts_path)
            if self.docling_artifacts_path is not None
            else None
        )

    async def parse(
        self,
        capability: str,
        body: bytes,
        *,
        target: str,
        limits: DocumentLimits,
        timeout_seconds: float,
    ) -> DocumentParseResult:
        if capability not in DOCUMENT_CAPABILITIES:
            raise AdapterExecutionError(
                f"document worker cannot execute {capability}"
            )
        if timeout_seconds <= 0:
            raise AdapterExecutionError("document parsing has no deadline budget")
        if len(body) > limits.maximum_input_bytes:
            raise AdapterExecutionError(
                "document body exceeds the worker byte limit"
            )
        expected_observed = _detect_capability(capability, target, body)
        if self.docling_artifacts_path is not None:
            try:
                current_root_id = docling_artifact_root_identity(
                    self.docling_artifacts_path
                )
            except (OSError, ValueError) as exc:
                raise AdapterExecutionError(
                    "Docling artifact bundle is unavailable or invalid"
                ) from exc
            if not hmac.compare_digest(
                current_root_id,
                self.docling_artifact_root_id or "",
            ):
                raise AdapterExecutionError(
                    "Docling artifact root identity changed"
                )
        docling = (
            {
                "artifact_bundle_id": self.docling_artifact_bundle_id,
                "artifacts_path": str(self.docling_artifacts_path),
                "document_timeout_seconds": min(timeout_seconds, 3600.0),
            }
            if (
                expected_observed in {"pdf", "scanned_pdf"}
                and self.docling_artifacts_path is not None
            )
            else None
        )
        payload = json.dumps(
            {
                "capability": capability,
                # The worker needs only a suffix hint. Hosts, path components,
                # query values, and opaque authentication material never cross
                # this boundary.
                "target_path": (
                    "file" + PurePosixPath(urlsplit(target).path).suffix
                ),
                "body": base64.b64encode(body).decode("ascii"),
                "limits": {
                    "maximum_input_bytes": limits.maximum_input_bytes,
                    "maximum_output_bytes": limits.maximum_output_bytes,
                    "maximum_blocks": limits.maximum_blocks,
                    "maximum_depth": limits.maximum_depth,
                    "maximum_archive_members": limits.maximum_archive_members,
                    "maximum_archive_ratio": limits.maximum_archive_ratio,
                },
                "docling": docling,
            },
            separators=(",", ":"),
            allow_nan=False,
        ).encode("utf-8")
        read_only_paths: tuple[Path, ...] = ()
        if docling is not None:
            assert self.docling_artifacts_path is not None
            read_only_paths = (self.docling_artifacts_path,)
        try:
            process = await run_bounded(
                (sys.executable, "-m", "fetech.document_worker"),
                payload,
                timeout_seconds=timeout_seconds,
                memory_mb=(
                    self.docling_memory_mb
                    if docling is not None
                    else self.memory_mb
                ),
                maximum_output_bytes=min(
                    100_000_000,
                    limits.maximum_output_bytes + 65_536,
                ),
                maximum_file_bytes=max(
                    limits.maximum_input_bytes,
                    min(
                        100_000_000,
                        limits.maximum_output_bytes + 65_536,
                    ),
                ),
                isolation=self.isolation.request(
                    WorkerIsolationProfile.DOCUMENT_PARSER,
                    read_only_paths=read_only_paths,
                ),
            )
        except LogicBackendError as exc:
            raise AdapterExecutionError(
                "bounded document parse process failed"
            ) from exc
        if process.returncode == 2:
            raise AdapterDependencyError(
                "document parsing requires the matching fetech[documents] dependency "
                "or configured connector"
            )
        if not process.stdout:
            raise AdapterExecutionError("document worker exited without output")
        try:
            response = json.loads(
                process.stdout,
                parse_constant=_reject_json_constant,
            )
        except (json.JSONDecodeError, ValueError) as exc:
            raise AdapterExecutionError(
                "document worker returned malformed output"
            ) from exc
        if process.returncode != 0:
            raise AdapterExecutionError("bounded document parsing failed")
        return _validate_worker_result(
            response,
            limits=limits,
            expected_observed=expected_observed,
            expected_docling_bundle_id=(
                self.docling_artifact_bundle_id
                if docling is not None
                else None
            ),
        )


class DocumentAdapter:
    """Normalize fetched document bytes without acquiring another resource."""

    def __init__(
        self,
        *,
        maximum_parse_bytes: int = 16_000_000,
        maximum_output_bytes: int = _MAX_DOCUMENT_OUTPUT_BYTES,
        maximum_blocks: int = 10_000,
        maximum_depth: int = 64,
        maximum_archive_members: int = 10_000,
        worker_timeout_seconds: float = 20.0,
        provider_timeout_seconds: float = 20.0,
        parser: DocumentParser | None = None,
        git_lfs_resolver: GitLFSResolver | None = None,
        pdf_ocr_provider: PDFOCRProvider | None = None,
    ) -> None:
        if min(
            maximum_parse_bytes,
            maximum_output_bytes,
            maximum_blocks,
            maximum_depth,
            maximum_archive_members,
            worker_timeout_seconds,
            provider_timeout_seconds,
        ) <= 0:
            raise ValueError("document adapter limits must be positive")
        self.maximum_parse_bytes = maximum_parse_bytes
        self.maximum_output_bytes = maximum_output_bytes
        self.maximum_blocks = maximum_blocks
        self.maximum_depth = maximum_depth
        self.maximum_archive_members = maximum_archive_members
        self.worker_timeout_seconds = worker_timeout_seconds
        self.provider_timeout_seconds = provider_timeout_seconds
        self.parser = parser or DocumentParseWorker()
        self.git_lfs_resolver = git_lfs_resolver
        self.pdf_ocr_provider = pdf_ocr_provider

    async def execute(self, node: PlanNode, context: ExecutionContext) -> None:
        attempt = FetchAttempt(
            capability_id=node.capability_id,
            adapter_version="0.4.0a0",
            sanitized_destination=sanitize_url_for_request(
                context.request.target,
                context.request,
            ),
            status=AttemptStatus.RUNNING,
        )
        attempt_index = len(context.attempts)
        context.attempts.append(attempt)
        try:
            if node.capability_id not in DOCUMENT_CAPABILITIES:
                raise AdapterExecutionError(
                    f"document adapter cannot execute {node.capability_id}"
                )
            raw = context.latest_artifact("raw")
            if raw is None or not context.resources:
                raise AdapterExecutionError(
                    "document parsing requires an acquired raw artifact"
                )
            source_resource = next(
                (
                    resource
                    for resource in reversed(context.resources)
                    if resource.resource_id == raw.source_resource_id
                ),
                None,
            )
            if source_resource is None:
                raise AdapterExecutionError(
                    "document source resource is missing from the execution context"
                )
            parse_limit = min(
                self.maximum_parse_bytes,
                _MAX_DOCUMENT_INPUT_BYTES,
                context.request.budget.bytes,
                context.request.budget.decompressed_bytes,
            )
            if raw.size > parse_limit:
                raise AdapterExecutionError(
                    "document body exceeds the parser byte limit"
                )
            body = await context.cas.get(raw.cas_uri, maximum_bytes=parse_limit)
            remaining_output_bytes = int(
                context.remaining_budget("decompressed_bytes")
            )
            context.require_budget("decompressed_bytes", 1)
            if node.capability_id == "github_raw":
                _validate_github_raw_target(context.request.target)
                _validate_github_raw_target(source_resource.canonical_url)
            git_lfs_resolved = False
            resolved_transfer_bytes = 0
            if node.capability_id == "git_lfs":
                pointer = _parse_git_lfs_pointer(body)
                if pointer is not None:
                    maximum_resolved_bytes = min(
                        parse_limit,
                        int(context.remaining_budget("bytes")),
                        remaining_output_bytes,
                    )
                    context.require_budget("bytes", 1)
                    body = await self._resolve_git_lfs(
                        pointer,
                        context=context,
                        maximum_bytes=maximum_resolved_bytes,
                        attempt_index=attempt_index,
                    )
                    git_lfs_resolved = True
                    resolved_transfer_bytes = len(body)
                    remaining_output_bytes -= resolved_transfer_bytes
                    context.require_budget("decompressed_bytes", 1)
            limits = DocumentLimits(
                maximum_input_bytes=parse_limit,
                maximum_output_bytes=min(
                    self.maximum_output_bytes,
                    remaining_output_bytes,
                ),
                maximum_blocks=min(
                    self.maximum_blocks,
                    max(1, context.request.budget.decompressed_bytes // 16),
                ),
                maximum_depth=self.maximum_depth,
                maximum_archive_members=min(
                    self.maximum_archive_members,
                    context.request.budget.archive_members,
                ),
                maximum_archive_ratio=context.request.budget.archive_ratio,
            )
            parsed = await self.parser.parse(
                node.capability_id,
                body,
                target=context.request.target,
                limits=limits,
                timeout_seconds=min(
                    self.worker_timeout_seconds,
                    context.request.budget.deadline_seconds,
                ),
            )
            parser_capability = parsed.observed_capability
            document = parsed.document
            locators = parsed.locators
            parser = parsed.parser
            parser_components = dict(parsed.parser_components)
            parser_identity = _document_parser_identity(
                parser,
                parser_components,
            )
            text = _document_text(document)
            quality = _assess_document(
                parser_capability,
                document,
                text,
                expected_language=context.request.language,
            )
            ocr_fallback = False
            if (
                quality.page_state == PageState.NEEDS_OCR
                and parser_capability in {"pdf", "scanned_pdf"}
                and self.pdf_ocr_provider is not None
            ):
                document, locators = await self._extract_pdf_ocr(
                    body,
                    page_count=len(document.get("blocks", [])),
                    language=context.request.language,
                    timeout_seconds=min(
                        self.provider_timeout_seconds,
                        context.request.budget.deadline_seconds,
                    ),
                    maximum_output_bytes=limits.maximum_output_bytes,
                )
                parser = "configured-pdf-ocr"
                parser_identity = "configured-pdf-ocr/0.4"
                text = _document_text(document)
                quality = _assess_document(
                    parser_capability,
                    document,
                    text,
                    expected_language=context.request.language,
                )
                ocr_fallback = True
                context.record_runtime_event(
                    "document_ocr_completed",
                    "documents",
                    capability_id=node.capability_id,
                    pages=len(locators),
                    accepted=quality.accepted,
                )
            if parsed.fallback_reason is not None:
                context.record_runtime_event(
                    "document_parser_fallback",
                    "documents",
                    capability_id=node.capability_id,
                    from_parser="docling",
                    to_parser=parser,
                    reason=parsed.fallback_reason,
                    artifact_bundle_id=parsed.artifact_bundle_id,
                )
            encoded = json.dumps(
                document,
                ensure_ascii=False,
                sort_keys=True,
                separators=(",", ":"),
                allow_nan=False,
            ).encode("utf-8")
            if len(encoded) > limits.maximum_output_bytes:
                raise AdapterExecutionError(
                    "normalized document exceeds the decompressed-byte budget"
                )
            context.require_budget("decompressed_bytes", len(encoded))
            uri, digest, size = await context.cas.put(encoded)
            artifact = build_artifact(
                role="primary" if quality.accepted else "checked-only",
                representation="document",
                media_type="application/vnd.fetech.document+json",
                cas_uri=uri,
                digest=digest,
                size=size,
                resource=source_resource,
                extractor=f"builtin-document/{parser_identity}",
                quality=quality,
                parents=(raw,),
                locators=locators,
            )
            context.artifacts.append(artifact)
            context.accepted = context.accepted or quality.accepted
            context.record_outcome(
                node.capability_id,
                CapabilityOutcomeStatus.APPLIED,
                "documents",
                observed_format=parser_capability,
                parser=parser,
                parser_identity=parser_identity,
                bounded=True,
                needs_ocr=quality.page_state == PageState.NEEDS_OCR,
                ocr_fallback=ocr_fallback,
                git_lfs_resolved=git_lfs_resolved,
                fallback_reason=parsed.fallback_reason,
                artifact_bundle_id=parsed.artifact_bundle_id,
                docling_version=parser_components.get("docling"),
                docling_slim_version=parser_components.get("docling-slim"),
                docling_core_version=parser_components.get("docling-core"),
                docling_parse_version=parser_components.get("docling-parse"),
                docling_ibm_models_version=parser_components.get(
                    "docling-ibm-models"
                ),
            )
            if parser_capability != node.capability_id:
                context.record_outcome(
                    parser_capability,
                    CapabilityOutcomeStatus.OBSERVED,
                    "documents",
                    requested_capability=node.capability_id,
                )
            context.record_quality_outcomes(quality, stage="quality")
            current_attempt = context.attempts[attempt_index]
            consumed_budget = dict(current_attempt.consumed_budget)
            warnings = list(current_attempt.warnings)
            if parsed.fallback_reason is not None:
                warnings.append(
                    f"document_parser_fallback:{parsed.fallback_reason}"
                )
            consumed_budget["decompressed_bytes"] = (
                consumed_budget.get("decompressed_bytes", 0) + len(encoded)
            )
            context.attempts[attempt_index] = current_attempt.model_copy(
                update={
                    "status": AttemptStatus.SUCCEEDED,
                    "finished_at": datetime.now(UTC),
                    "bytes_received": len(body),
                    "parser": parser,
                    "artifact_ids": (artifact.artifact_id,),
                    "warnings": tuple(warnings),
                    "consumed_budget": consumed_budget,
                }
            )
        except (ImportError, AdapterDependencyError) as exc:
            context.record_outcome(
                node.capability_id,
                CapabilityOutcomeStatus.DEPENDENCY_MISSING,
                "documents",
                reason=str(exc),
            )
            current_attempt = context.attempts[attempt_index]
            context.attempts[attempt_index] = current_attempt.model_copy(
                update={
                    "status": AttemptStatus.FAILED,
                    "finished_at": datetime.now(UTC),
                    "failure_code": "dependency_missing",
                    "warnings": (str(exc),),
                }
            )
            if isinstance(exc, AdapterDependencyError):
                raise
            raise AdapterDependencyError(str(exc)) from exc
        except (OSError, ValueError, KeyError, ElementTree.ParseError) as exc:
            context.record_outcome(
                node.capability_id,
                CapabilityOutcomeStatus.FAILED,
                "documents",
                reason=str(exc),
            )
            current_attempt = context.attempts[attempt_index]
            context.attempts[attempt_index] = current_attempt.model_copy(
                update={
                    "status": AttemptStatus.FAILED,
                    "finished_at": datetime.now(UTC),
                    "failure_code": "malformed_document",
                    "warnings": (str(exc),),
                }
            )
            raise AdapterExecutionError(str(exc)) from exc
        except AdapterExecutionError as exc:
            context.record_outcome(
                node.capability_id,
                CapabilityOutcomeStatus.FAILED,
                "documents",
                reason=str(exc),
            )
            current_attempt = context.attempts[attempt_index]
            context.attempts[attempt_index] = current_attempt.model_copy(
                update={
                    "status": AttemptStatus.FAILED,
                    "finished_at": datetime.now(UTC),
                    "failure_code": "document_error",
                    "warnings": (str(exc),),
                }
            )
            raise

    async def _resolve_git_lfs(
        self,
        pointer: GitLFSPointer,
        *,
        context: ExecutionContext,
        maximum_bytes: int,
        attempt_index: int,
    ) -> bytes:
        resolver = self.git_lfs_resolver
        if resolver is None:
            raise AdapterDependencyError(
                "Git LFS pointer resolution requires a configured exact-origin resolver"
            )
        if pointer.size > maximum_bytes:
            raise AdapterExecutionError(
                "Git LFS declared size exceeds the request byte budget"
            )
        try:
            origin = canonical_origin(context.request.target)
        except ValueError as exc:
            raise AdapterExecutionError("Git LFS target origin is invalid") from exc
        sanitized_target = sanitize_url_for_request(
            context.request.target,
            context.request,
        )
        if (
            len(sanitized_target.encode("utf-8")) > _MAX_PROVIDER_TARGET_BYTES
            or canonical_origin(sanitized_target) != origin
        ):
            raise AdapterExecutionError("Git LFS sanitized target is invalid")
        timeout_seconds = min(
            self.provider_timeout_seconds,
            context.request.budget.deadline_seconds,
        )
        request = GitLFSResolveRequest(
            sanitized_target=sanitized_target,
            origin=origin,
            pointer=pointer,
            maximum_bytes=maximum_bytes,
            timeout_seconds=timeout_seconds,
        )
        try:
            resolved = await asyncio.wait_for(
                resolver.resolve(request),
                timeout=timeout_seconds,
            )
        except TimeoutError as exc:
            raise AdapterExecutionError("configured Git LFS resolver timed out") from exc
        except Exception as exc:
            raise AdapterExecutionError("configured Git LFS resolver failed") from exc
        if not isinstance(resolved, GitLFSResolvedObject):
            raise AdapterExecutionError(
                "configured Git LFS resolver returned invalid output"
            )
        if not isinstance(resolved.origin, str) or not isinstance(resolved.body, bytes):
            raise AdapterExecutionError(
                "configured Git LFS resolver returned invalid output"
            )
        context.record_attempt_consumption(
            attempt_index,
            bytes=len(resolved.body),
            decompressed_bytes=len(resolved.body),
        )
        if len(resolved.body) > maximum_bytes:
            raise AdapterBudgetExceededError(
                "configured Git LFS resolver exceeded its body byte budget"
            )
        try:
            resolved_origin = canonical_origin(resolved.origin, origin_only=True)
        except ValueError as exc:
            raise AdapterExecutionError(
                "configured Git LFS resolver returned invalid output"
            ) from exc
        if resolved_origin != origin:
            raise AdapterExecutionError(
                "configured Git LFS resolver violated the exact-origin contract"
            )
        if len(resolved.body) > maximum_bytes or len(resolved.body) != pointer.size:
            raise AdapterExecutionError(
                "configured Git LFS resolver returned a size mismatch"
            )
        digest = sha256(resolved.body).hexdigest()
        if not hmac.compare_digest(digest, pointer.oid_sha256):
            raise AdapterExecutionError(
                "configured Git LFS resolver returned a hash mismatch"
            )
        return resolved.body

    async def _extract_pdf_ocr(
        self,
        body: bytes,
        *,
        page_count: int,
        language: str | None,
        timeout_seconds: float,
        maximum_output_bytes: int,
    ) -> tuple[dict[str, Any], tuple[str, ...]]:
        provider = self.pdf_ocr_provider
        if provider is None:
            raise AdapterDependencyError("a configured PDF OCR provider is required")
        try:
            pages = await asyncio.wait_for(
                provider.extract_pdf(
                    body,
                    page_count=page_count,
                    language=language,
                    timeout_seconds=timeout_seconds,
                    maximum_output_bytes=maximum_output_bytes,
                ),
                timeout=timeout_seconds,
            )
        except TimeoutError as exc:
            raise AdapterExecutionError("configured PDF OCR provider timed out") from exc
        except Exception as exc:
            raise AdapterExecutionError("configured PDF OCR provider failed") from exc
        try:
            return _validate_pdf_ocr_pages(
                pages,
                page_count=page_count,
                maximum_blocks=self.maximum_blocks,
                maximum_output_bytes=maximum_output_bytes,
            )
        except AdapterExecutionError:
            raise
        except Exception as exc:
            raise AdapterExecutionError(
                "configured PDF OCR provider returned invalid output"
            ) from exc


def _parse_git_lfs_pointer(body: bytes) -> GitLFSPointer | None:
    if not body.startswith(_GIT_LFS_POINTER_PREFIX):
        return None
    if len(body) > _MAX_GIT_LFS_POINTER_BYTES:
        raise AdapterExecutionError("Git LFS pointer exceeds its bounded grammar")
    match = _GIT_LFS_POINTER.fullmatch(body)
    if match is None:
        raise AdapterExecutionError("Git LFS pointer grammar is invalid")
    return GitLFSPointer(
        oid_sha256=match.group("oid").decode("ascii"),
        size=int(match.group("size")),
    )


def _validate_github_raw_target(target: str) -> None:
    try:
        parts = urlsplit(target)
        origin = canonical_origin(target)
        decoded_path = unquote_to_bytes(parts.path).decode("utf-8")
    except (UnicodeDecodeError, ValueError) as exc:
        raise AdapterExecutionError("GitHub raw target is invalid") from exc
    if (
        origin != _GITHUB_RAW_ORIGIN
        or parts.query
        or parts.fragment
        or len(target.encode("utf-8")) > _MAX_PROVIDER_TARGET_BYTES
        or re.search(r"%(?![0-9A-Fa-f]{2})", parts.path)
        or re.search(r"%(?:2f|5c)", parts.path, re.IGNORECASE)
    ):
        raise AdapterExecutionError(
            "GitHub raw target must use the exact raw.githubusercontent.com origin"
        )
    components = decoded_path.removeprefix("/").split("/")
    if (
        len(components) < 4
        or any(
            not component
            or component in {".", ".."}
            or len(component.encode("utf-8")) > 1_024
            or "\\" in component
            or "\x00" in component
            or any(ord(character) < 32 or ord(character) == 127 for character in component)
            for component in components
        )
        or _GITHUB_PATH_COMPONENT.fullmatch(components[0]) is None
        or _GITHUB_PATH_COMPONENT.fullmatch(components[1]) is None
    ):
        raise AdapterExecutionError(
            "GitHub raw target path must identify an owner, repository, ref, and file"
        )


def _validate_pdf_ocr_pages(
    value: object,
    *,
    page_count: int,
    maximum_blocks: int,
    maximum_output_bytes: int,
) -> tuple[dict[str, Any], tuple[str, ...]]:
    if (
        not isinstance(value, Sequence)
        or isinstance(value, str | bytes | bytearray)
        or page_count <= 0
        or len(value) == 0
        or len(value) > min(page_count, maximum_blocks)
    ):
        raise AdapterExecutionError(
            "configured PDF OCR provider returned invalid output"
        )
    blocks: list[dict[str, Any]] = []
    seen_pages: set[int] = set()
    text_bytes = 0
    for item in value:
        if (
            not isinstance(item, PDFOCRPage)
            or not isinstance(item.locator, str)
            or not isinstance(item.text, str)
            or not item.text.strip()
        ):
            raise AdapterExecutionError(
                "configured PDF OCR provider returned invalid output"
            )
        locator_match = _PDF_PAGE_LOCATOR.fullmatch(item.locator)
        if locator_match is None:
            raise AdapterExecutionError(
                "configured PDF OCR provider returned invalid page locators"
            )
        page = int(locator_match.group("page"))
        if page > page_count or page in seen_pages:
            raise AdapterExecutionError(
                "configured PDF OCR provider returned invalid page locators"
            )
        seen_pages.add(page)
        text_bytes += len(item.text.encode("utf-8"))
        if text_bytes > maximum_output_bytes:
            raise AdapterExecutionError(
                "configured PDF OCR provider exceeded its output byte budget"
            )
        blocks.append({"locator": item.locator, "text": item.text})
    blocks.sort(key=lambda block: int(str(block["locator"]).split(":", 1)[1]))
    locators = _locators(blocks)
    document: dict[str, Any] = {"type": "pdf", "blocks": blocks}
    try:
        encoded = json.dumps(
            document,
            ensure_ascii=False,
            sort_keys=True,
            separators=(",", ":"),
            allow_nan=False,
        ).encode("utf-8")
    except (TypeError, ValueError) as exc:
        raise AdapterExecutionError(
            "configured PDF OCR provider returned invalid output"
        ) from exc
    if len(encoded) > maximum_output_bytes:
        raise AdapterExecutionError(
            "configured PDF OCR provider exceeded its output byte budget"
        )
    return document, locators


def _detect_capability(requested: str, target: str, body: bytes) -> str:
    """Resolve aggregate/connector IDs using bytes before filename hints."""

    if requested not in {"dataset_file", "github_raw", "git_lfs", "document_router"}:
        return requested
    if requested == "git_lfs" and _GIT_LFS_POINTER.fullmatch(body):
        return "git_lfs"
    return _detect_document_format(target, body)


def _detect_document_format(target: str, body: bytes) -> str:
    suffix = PurePosixPath(urlsplit(target).path.lower()).suffix
    if body.startswith(b"%PDF-"):
        return "pdf"
    # Classification remains a lightweight magic-byte probe in the daemon.
    # The hostile ZIP central directory is parsed only inside the worker.
    if body.startswith(_ZIP_SIGNATURES):
        return {
            ".docx": "docx",
            ".pptx": "pptx",
            ".xlsx": "xlsx",
        }.get(suffix, "zip_archive")
    if suffix in _MARKDOWN_SUFFIXES:
        return "markdown"
    if suffix in _TEXT_SUFFIXES:
        return "txt"
    by_suffix = {
        ".csv": "csv",
        ".docx": "docx",
        ".json": "json_file",
        ".pdf": "pdf",
        ".pptx": "pptx",
        ".xlsx": "xlsx",
        ".xml": "xml_file",
        ".zip": "zip_archive",
    }.get(suffix)
    if by_suffix is not None:
        return by_suffix
    prefix = body[:4_096]
    if _JSON_XML_PREFIX.search(prefix):
        stripped = prefix.lstrip()
        return "xml_file" if stripped.startswith(b"<") else "json_file"
    if _looks_like_csv(body):
        return "csv"
    return "txt" if _looks_like_text(body) else "dataset_file"


def _parse(
    capability: str,
    body: bytes,
    *,
    target: str = "",
    limits: DocumentLimits | None = None,
) -> tuple[dict[str, Any], tuple[str, ...], str]:
    limits = limits or DocumentLimits()
    if len(body) > limits.maximum_input_bytes:
        raise ValueError("document body exceeds the parser byte limit")
    if capability in {"dataset_file", "github_raw"}:
        capability = _detect_document_format(target, body)
    elif capability == "git_lfs":
        if _GIT_LFS_POINTER.fullmatch(body):
            raise ImportError(
                "Git LFS pointer resolution requires a configured origin-scoped connector"
            )
        capability = _detect_document_format(target, body)

    if capability in {"plain_text_file", "txt", "markdown"}:
        return _parse_text(capability, body, limits)
    if capability == "dataset_file":
        return _parse_binary_dataset(body)
    if capability == "csv":
        return _parse_csv(body, limits)
    if capability == "json_file":
        return _parse_json(body, limits)
    if capability == "xml_file":
        return _parse_xml(body, limits)
    if capability in {"pdf", "scanned_pdf"}:
        return _parse_pdf(body, limits)
    if capability == "docx":
        return _parse_docx(body, limits)
    if capability == "pptx":
        return _parse_pptx(body, limits)
    if capability == "xlsx":
        return _parse_xlsx(body, limits)
    if capability == "zip_archive":
        return _parse_zip(body, limits)
    raise ValueError(f"unsupported document capability: {capability}")


def _validate_worker_result(
    response: object,
    *,
    limits: DocumentLimits,
    expected_observed: str,
    expected_docling_bundle_id: str | None = None,
) -> DocumentParseResult:
    if not isinstance(response, dict):
        raise AdapterExecutionError("document worker response must be an object")
    raw_document = response.get("document")
    raw_locators = response.get("locators")
    parser = response.get("parser")
    observed = response.get("observed_capability")
    raw_parser_components = response.get("parser_components", {})
    artifact_bundle_id = response.get("artifact_bundle_id")
    fallback_reason = response.get("fallback_reason")
    if not isinstance(raw_document, dict):
        raise AdapterExecutionError("document worker omitted its document")
    if set(raw_document) != {"blocks", "type"}:
        raise AdapterExecutionError("document worker document schema is invalid")
    blocks = raw_document.get("blocks")
    if not isinstance(blocks, list) or len(blocks) > limits.maximum_blocks:
        raise AdapterExecutionError("document worker returned invalid blocks")
    block_locators: list[str] = []
    for block in blocks:
        if not isinstance(block, dict) or not isinstance(block.get("locator"), str):
            raise AdapterExecutionError(
                "document worker returned a block without a locator"
            )
        block_locators.append(block["locator"])
    if (
        not isinstance(raw_locators, list)
        or not all(isinstance(locator, str) for locator in raw_locators)
        or raw_locators != block_locators
        or len(raw_locators) != len(set(raw_locators))
    ):
        raise AdapterExecutionError("document worker locators are invalid")
    if not isinstance(observed, str) or observed != expected_observed:
        raise AdapterExecutionError(
            "document worker observed format contradicts deterministic routing"
        )
    try:
        expected_parser, expected_type, locator_prefix = _DOCUMENT_RESULT_CONTRACTS[
            observed
        ]
    except KeyError as exc:
        raise AdapterExecutionError(
            "document worker observed format is not a normalized document"
        ) from exc
    if parser != expected_parser and not (
        observed in {"pdf", "scanned_pdf"}
        and parser == "docling"
        and expected_docling_bundle_id is not None
    ):
        raise AdapterExecutionError("document worker parser identity is invalid")
    parser_components = _validate_parser_provenance(
        parser=parser,
        observed=observed,
        raw_components=raw_parser_components,
        artifact_bundle_id=artifact_bundle_id,
        fallback_reason=fallback_reason,
        expected_docling_bundle_id=expected_docling_bundle_id,
    )
    if raw_document.get("type") != expected_type:
        raise AdapterExecutionError("document worker document type is invalid")
    if any(not locator.startswith(locator_prefix) for locator in raw_locators):
        raise AdapterExecutionError("document worker locator family is invalid")
    if observed in {"pdf", "scanned_pdf"} and any(
        (match := _PDF_PAGE_LOCATOR.fullmatch(locator)) is None
        or int(match.group("page")) > limits.maximum_blocks
        for locator in raw_locators
    ):
        raise AdapterExecutionError("document worker page locators are invalid")
    for block in blocks:
        _validate_worker_block(block, observed=observed)
    try:
        encoded = json.dumps(
            raw_document,
            ensure_ascii=False,
            sort_keys=True,
            separators=(",", ":"),
            allow_nan=False,
        ).encode("utf-8")
    except (TypeError, ValueError) as exc:
        raise AdapterExecutionError(
            "document worker returned a non-serializable document"
        ) from exc
    if len(encoded) > limits.maximum_output_bytes:
        raise AdapterExecutionError(
            "document worker exceeded the decompressed-byte budget"
        )
    return DocumentParseResult(
        document=raw_document,
        locators=tuple(raw_locators),
        parser=parser,
        observed_capability=observed,
        parser_components=parser_components,
        artifact_bundle_id=artifact_bundle_id,
        fallback_reason=fallback_reason,
    )


def _validate_docling_artifacts_path(value: Path | None) -> Path | None:
    if value is None:
        return None
    if not isinstance(value, Path):
        raise TypeError("Docling artifacts path must be a pathlib.Path")
    supplied = value.expanduser()
    try:
        if supplied.is_symlink():
            raise ValueError("Docling artifacts path cannot be a symbolic link")
        resolved = supplied.resolve(strict=True)
    except OSError as exc:
        raise ValueError("Docling artifacts path must be an existing directory") from exc
    if (
        not resolved.is_dir()
        or resolved == Path(resolved.anchor)
        or len(str(resolved).encode("utf-8")) > _MAX_DOCLING_ARTIFACTS_PATH_BYTES
        or "\x00" in str(resolved)
    ):
        raise ValueError(
            "Docling artifacts path must be a bounded non-root directory"
        )
    return resolved


def _docling_artifact_bundle_identity(path: Path) -> str:
    """Return the portable SHA-256 identity of the complete model contents."""

    return inspect_docling_artifact_bundle(
        path,
        require_manifest=True,
    ).bundle_sha256


def _validate_parser_provenance(
    *,
    parser: object,
    observed: object,
    raw_components: object,
    artifact_bundle_id: object,
    fallback_reason: object,
    expected_docling_bundle_id: str | None,
) -> tuple[tuple[str, str], ...]:
    if not isinstance(raw_components, dict):
        raise AdapterExecutionError("document worker parser provenance is invalid")
    if (
        not all(isinstance(key, str) and isinstance(value, str) for key, value in raw_components.items())
        or not set(raw_components).issubset(_DOCLING_COMPONENT_KEYS)
        or any(
            _SAFE_COMPONENT_VERSION.fullmatch(value) is None
            for value in raw_components.values()
        )
    ):
        raise AdapterExecutionError("document worker parser provenance is invalid")
    if artifact_bundle_id is not None and (
        not isinstance(artifact_bundle_id, str)
        or _SHA256_HEX.fullmatch(artifact_bundle_id) is None
    ):
        raise AdapterExecutionError("document worker artifact identity is invalid")
    if fallback_reason is not None and (
        not isinstance(fallback_reason, str)
        or fallback_reason not in _DOCLING_FALLBACK_REASONS
    ):
        raise AdapterExecutionError("document worker fallback reason is invalid")

    if parser == "docling":
        if (
            observed not in {"pdf", "scanned_pdf"}
            or not _DOCLING_REQUIRED_COMPONENT_KEYS.issubset(raw_components)
            or any(
                raw_components.get(name) != expected_version
                for name, expected_version in (
                    _DOCLING_REQUIRED_COMPONENT_VERSIONS.items()
                )
            )
            or artifact_bundle_id is None
            or expected_docling_bundle_id is None
            or not hmac.compare_digest(
                artifact_bundle_id,
                expected_docling_bundle_id,
            )
            or fallback_reason is not None
        ):
            raise AdapterExecutionError("document worker Docling provenance is invalid")
    elif raw_components:
        raise AdapterExecutionError("document worker returned foreign parser provenance")

    if fallback_reason is not None and (
        parser != "pypdf"
        or observed not in {"pdf", "scanned_pdf"}
        or artifact_bundle_id is None
        or expected_docling_bundle_id is None
        or not hmac.compare_digest(
            artifact_bundle_id,
            expected_docling_bundle_id,
        )
    ):
        raise AdapterExecutionError("document worker fallback provenance is invalid")
    if artifact_bundle_id is not None and parser != "docling" and fallback_reason is None:
        raise AdapterExecutionError("document worker artifact identity is unexpected")
    if expected_docling_bundle_id is None and (
        artifact_bundle_id is not None or fallback_reason is not None
    ):
        raise AdapterExecutionError(
            "document worker returned unconfigured Docling provenance"
        )
    if (
        expected_docling_bundle_id is not None
        and parser != "docling"
        and fallback_reason is None
    ):
        raise AdapterExecutionError(
            "document worker omitted configured Docling provenance"
        )
    return tuple(sorted(raw_components.items()))


def _document_parser_identity(
    parser: str,
    parser_components: dict[str, str],
) -> str:
    if parser == "docling":
        return f"docling-slim/{parser_components['docling-slim']}"
    return f"{parser}/0.4"


def _validate_worker_block(block: dict[str, object], *, observed: str) -> None:
    text_formats = {"docx", "markdown", "pdf", "pptx", "scanned_pdf", "txt"}
    if observed in text_formats:
        if set(block) != {"locator", "text"} or not isinstance(block["text"], str):
            raise AdapterExecutionError("document worker text block is invalid")
        return
    if observed in {"csv", "xlsx"}:
        cells = block.get("cells")
        if set(block) != {"cells", "locator"} or not isinstance(cells, list):
            raise AdapterExecutionError("document worker table block is invalid")
        if not all(_is_json_scalar(value) for value in cells):
            raise AdapterExecutionError(
                "document worker table cell is not a finite scalar"
            )
        return
    if observed == "json_file":
        if set(block) != {"locator", "value"}:
            raise AdapterExecutionError("document worker JSON block is invalid")
        value = block["value"]
        if not (
            _is_json_scalar(value)
            or value == {}
            or value == []
        ):
            raise AdapterExecutionError(
                "document worker JSON leaf is not a bounded value"
            )
        return
    if observed == "xml_file":
        if not {"locator", "tag"}.issubset(block) or not set(block).issubset(
            {"attributes", "locator", "tag", "text"}
        ):
            raise AdapterExecutionError("document worker XML block is invalid")
        if not isinstance(block["tag"], str):
            raise AdapterExecutionError("document worker XML tag is invalid")
        if "text" in block and not isinstance(block["text"], str):
            raise AdapterExecutionError("document worker XML text is invalid")
        attributes = block.get("attributes", {})
        if not isinstance(attributes, dict) or not all(
            isinstance(key, str) and isinstance(value, str)
            for key, value in attributes.items()
        ):
            raise AdapterExecutionError("document worker XML attributes are invalid")
        return
    if observed in {"dataset_file", "zip_archive"}:
        size = block.get("size")
        digest = block.get("sha256")
        if (
            set(block) != {"locator", "sha256", "size"}
            or not isinstance(size, int)
            or isinstance(size, bool)
            or size < 0
            or not isinstance(digest, str)
            or re.fullmatch(r"[0-9a-f]{64}", digest) is None
        ):
            raise AdapterExecutionError("document worker binary block is invalid")
        return
    raise AdapterExecutionError("document worker block format is unsupported")


def _is_json_scalar(value: object) -> bool:
    if value is None or isinstance(value, (str, bool, int)):
        return True
    return isinstance(value, float) and math.isfinite(value)


def _parse_text(
    capability: str,
    body: bytes,
    limits: DocumentLimits,
) -> tuple[dict[str, Any], tuple[str, ...], str]:
    text = body.decode("utf-8-sig", errors="replace")
    lines = text.splitlines(keepends=False)
    if not lines:
        lines = [""]
    _enforce_block_limit(len(lines), limits)
    blocks = [
        {"locator": f"line:{index}", "text": line}
        for index, line in enumerate(lines, start=1)
    ]
    document_type = "markdown" if capability == "markdown" else "text"
    parser = "markdown-text" if capability == "markdown" else "text"
    return {"type": document_type, "blocks": blocks}, _locators(blocks), parser


def _parse_binary_dataset(
    body: bytes,
) -> tuple[dict[str, Any], tuple[str, ...], str]:
    locator = f"byte:0-{len(body)}"
    blocks = [
        {
            "locator": locator,
            "size": len(body),
            "sha256": sha256(body).hexdigest(),
        }
    ]
    return (
        {"type": "binary_dataset", "blocks": blocks},
        (locator,),
        "binary-metadata",
    )


def _parse_csv(
    body: bytes,
    limits: DocumentLimits,
) -> tuple[dict[str, Any], tuple[str, ...], str]:
    rows = csv.reader(io.StringIO(body.decode("utf-8-sig", errors="replace")))
    blocks: list[dict[str, Any]] = []
    for index, row in enumerate(rows, start=1):
        if index > limits.maximum_blocks:
            raise ValueError("document block limit exceeded")
        blocks.append(
            {
                "locator": f"row:{index}",
                "cells": [_json_scalar(value) for value in row],
            }
        )
    return {"type": "table", "blocks": blocks}, _locators(blocks), "csv"


def _parse_json(
    body: bytes,
    limits: DocumentLimits,
) -> tuple[dict[str, Any], tuple[str, ...], str]:
    try:
        payload = json.loads(
            body.decode("utf-8-sig"),
            object_pairs_hook=_unique_json_object,
            parse_constant=_reject_json_constant,
        )
    except (UnicodeDecodeError, RecursionError) as exc:
        raise ValueError("JSON input is not bounded valid UTF-8") from exc
    blocks: list[dict[str, Any]] = []
    stack: list[tuple[str, object, int]] = [("", payload, 0)]
    visited = 0
    while stack:
        pointer, value, depth = stack.pop()
        if depth > limits.maximum_depth:
            raise ValueError("JSON nesting depth limit exceeded")
        visited += 1
        if visited > limits.maximum_blocks:
            raise ValueError("document block limit exceeded")
        locator = f"json:{pointer or '/'}"
        if isinstance(value, dict):
            if not value:
                blocks.append({"locator": locator, "value": {}})
            for key in reversed(sorted(value)):
                child = f"{pointer}/{_json_pointer_token(key)}"
                stack.append((child, value[key], depth + 1))
        elif isinstance(value, list):
            if not value:
                blocks.append({"locator": locator, "value": []})
            for index in range(len(value) - 1, -1, -1):
                stack.append((f"{pointer}/{index}", value[index], depth + 1))
        else:
            blocks.append({"locator": locator, "value": value})
    return (
        {"type": "json", "blocks": blocks},
        _locators(blocks),
        "stdlib-json",
    )


def _parse_xml(
    body: bytes,
    limits: DocumentLimits,
) -> tuple[dict[str, Any], tuple[str, ...], str]:
    if _FORBIDDEN_XML.search(body):
        raise ValueError("XML DTD and entity declarations are forbidden")
    try:
        root = ElementTree.fromstring(body)
    except ElementTree.ParseError:
        raise
    blocks: list[dict[str, Any]] = []
    stack: list[tuple[ElementTree.Element, str, int]] = [
        (root, f"/{root.tag}[1]", 0)
    ]
    while stack:
        element, path, depth = stack.pop()
        if depth > limits.maximum_depth:
            raise ValueError("XML nesting depth limit exceeded")
        if len(blocks) >= limits.maximum_blocks:
            raise ValueError("document block limit exceeded")
        block: dict[str, Any] = {
            "locator": f"xpath:{path}",
            "tag": element.tag,
        }
        text = (element.text or "").strip()
        if text:
            block["text"] = text
        if element.attrib:
            block["attributes"] = {
                key: element.attrib[key] for key in sorted(element.attrib)
            }
        blocks.append(block)
        counts: dict[str, int] = {}
        children: list[tuple[ElementTree.Element, str, int]] = []
        for child in element:
            counts[child.tag] = counts.get(child.tag, 0) + 1
            child_path = f"{path}/{child.tag}[{counts[child.tag]}]"
            children.append((child, child_path, depth + 1))
        stack.extend(reversed(children))
    return {"type": "xml", "blocks": blocks}, _locators(blocks), "stdlib-xml"


def _parse_pdf(
    body: bytes,
    limits: DocumentLimits,
) -> tuple[dict[str, Any], tuple[str, ...], str]:
    try:
        from pypdf import PdfReader
    except ImportError as exc:
        raise ImportError("install fetech[documents] for PDF parsing") from exc
    reader = PdfReader(io.BytesIO(body), strict=True)
    if len(reader.pages) > limits.maximum_blocks:
        raise ValueError("document block limit exceeded")
    blocks: list[dict[str, Any]] = []
    remaining = limits.maximum_output_bytes
    for index, page in enumerate(reader.pages, start=1):
        text = page.extract_text() or ""
        encoded_size = len(text.encode("utf-8"))
        if encoded_size > remaining:
            raise ValueError("PDF extracted text exceeds the decompressed-byte budget")
        remaining -= encoded_size
        blocks.append({"locator": f"page:{index}", "text": text})
    return {"type": "pdf", "blocks": blocks}, _locators(blocks), "pypdf"


def _parse_docx(
    body: bytes,
    limits: DocumentLimits,
) -> tuple[dict[str, Any], tuple[str, ...], str]:
    _validate_zip_container(body, limits, expected_suffix=".docx")
    try:
        from docx import Document
    except ImportError as exc:
        raise ImportError("install fetech[documents] for DOCX parsing") from exc
    paragraphs = Document(io.BytesIO(body)).paragraphs
    if len(paragraphs) > limits.maximum_blocks:
        raise ValueError("document block limit exceeded")
    blocks = [
        {"locator": f"paragraph:{index}", "text": paragraph.text}
        for index, paragraph in enumerate(paragraphs, start=1)
        if paragraph.text.strip()
    ]
    return {"type": "docx", "blocks": blocks}, _locators(blocks), "python-docx"


def _parse_pptx(
    body: bytes,
    limits: DocumentLimits,
) -> tuple[dict[str, Any], tuple[str, ...], str]:
    _validate_zip_container(body, limits, expected_suffix=".pptx")
    try:
        from pptx import Presentation
    except ImportError as exc:
        raise ImportError("install fetech[documents] for PPTX parsing") from exc
    slides = Presentation(io.BytesIO(body)).slides
    if len(slides) > limits.maximum_blocks:
        raise ValueError("document block limit exceeded")
    blocks: list[dict[str, Any]] = []
    for slide_index, slide in enumerate(slides, start=1):
        text = "\n".join(
            shape.text
            for shape in slide.shapes
            if hasattr(shape, "text") and shape.text.strip()
        )
        blocks.append({"locator": f"slide:{slide_index}", "text": text})
    return {"type": "pptx", "blocks": blocks}, _locators(blocks), "python-pptx"


def _parse_xlsx(
    body: bytes,
    limits: DocumentLimits,
) -> tuple[dict[str, Any], tuple[str, ...], str]:
    _validate_zip_container(body, limits, expected_suffix=".xlsx")
    try:
        from openpyxl import load_workbook
    except ImportError as exc:
        raise ImportError("install fetech[documents] for XLSX parsing") from exc
    workbook = load_workbook(
        io.BytesIO(body),
        read_only=True,
        data_only=True,
        keep_links=False,
    )
    blocks: list[dict[str, Any]] = []
    try:
        for sheet in workbook.worksheets:
            for row_index, row in enumerate(sheet.iter_rows(values_only=True), start=1):
                if len(blocks) >= limits.maximum_blocks:
                    raise ValueError("document block limit exceeded")
                blocks.append(
                    {
                        "locator": f"sheet:{sheet.title}/row:{row_index}",
                        "cells": [_json_scalar(value) for value in row],
                    }
                )
    finally:
        workbook.close()
    return {"type": "xlsx", "blocks": blocks}, _locators(blocks), "openpyxl"


def _parse_zip(
    body: bytes,
    limits: DocumentLimits,
) -> tuple[dict[str, Any], tuple[str, ...], str]:
    infos = _validate_zip_container(body, limits)
    blocks: list[dict[str, Any]] = []
    remaining = limits.maximum_output_bytes
    with zipfile.ZipFile(io.BytesIO(body)) as archive:
        for info in infos:
            if info.is_dir():
                continue
            digest, actual_size = _zip_member_digest(
                archive,
                info,
                maximum_bytes=remaining,
            )
            remaining -= actual_size
            blocks.append(
                {
                    "locator": f"member:{info.filename}",
                    "size": actual_size,
                    "sha256": digest,
                }
            )
    return (
        {"type": "zip_archive", "blocks": blocks},
        _locators(blocks),
        "stdlib-zip",
    )


def _validate_zip_container(
    body: bytes,
    limits: DocumentLimits,
    *,
    expected_suffix: str | None = None,
) -> list[zipfile.ZipInfo]:
    if not zipfile.is_zipfile(io.BytesIO(body)):
        label = expected_suffix or ".zip"
        raise ValueError(f"malformed {label} ZIP container")
    expanded = 0
    with zipfile.ZipFile(io.BytesIO(body)) as archive:
        infos = archive.infolist()
        if len(infos) > limits.maximum_archive_members:
            raise ValueError("archive member limit exceeded")
        names = [info.filename for info in infos]
        if len(names) != len(set(names)):
            raise ValueError("duplicate archive member names are forbidden")
        for info in infos:
            _validate_member_name(info.filename)
            mode = info.external_attr >> 16
            if stat.S_ISLNK(mode):
                raise ValueError("archive symlinks are forbidden")
            if info.flag_bits & 0x1:
                raise ValueError("encrypted archive members are unsupported")
            if info.is_dir():
                continue
            ratio = info.file_size / max(1, info.compress_size)
            if ratio > limits.maximum_archive_ratio:
                raise ValueError("archive compression ratio exceeded")
            expanded += info.file_size
            if expanded > limits.maximum_output_bytes:
                raise ValueError("archive expanded-byte limit exceeded")
        if expected_suffix is not None:
            _validate_ooxml_identity(archive, expected_suffix)
        return infos


def _validate_ooxml_identity(
    archive: zipfile.ZipFile,
    expected_suffix: str,
) -> None:
    required = {
        ".docx": "word/document.xml",
        ".pptx": "ppt/presentation.xml",
        ".xlsx": "xl/workbook.xml",
    }[expected_suffix]
    if required not in archive.namelist():
        raise ValueError(f"ZIP container is not a valid {expected_suffix} document")


def _zip_member_digest(
    archive: zipfile.ZipFile,
    info: zipfile.ZipInfo,
    *,
    maximum_bytes: int,
) -> tuple[str, int]:
    with archive.open(info, "r") as member:
        digest = sha256()
        size = 0
        while chunk := member.read(64 * 1024):
            size += len(chunk)
            if size > maximum_bytes or size > info.file_size:
                raise ValueError("archive expanded-byte limit exceeded")
            digest.update(chunk)
        if size != info.file_size:
            raise ValueError("archive member size contradicts its directory entry")
        return digest.hexdigest(), size


def _assess_document(
    capability: str,
    document: dict[str, Any],
    text: str,
    *,
    expected_language: str | None,
) -> QualityAssessment:
    if capability in {"pdf", "scanned_pdf"} and not text.strip():
        return QualityAssessment(
            page_state=PageState.NEEDS_OCR,
            score=0,
            accepted=False,
            completeness=0,
            reasons=("PDF contains no extractable text; OCR is required",),
        )
    if capability in {"dataset_file", "json_file", "xml_file", "zip_archive"}:
        blocks = document.get("blocks", [])
        accepted = bool(blocks)
        if capability == "dataset_file" and accepted:
            accepted = bool(blocks[0].get("size"))
        return QualityAssessment(
            page_state=PageState.OK if accepted else PageState.EMPTY,
            score=1.0 if accepted else 0.0,
            accepted=accepted,
            completeness=1.0 if accepted else 0.0,
            reasons=() if accepted else ("document contains no records",),
        )
    return assess_text(text, expected_language=expected_language)


def _document_text(document: dict[str, Any]) -> str:
    parts: list[str] = []
    for block in document.get("blocks", []):
        if "text" in block:
            parts.append(str(block["text"]))
        elif "cells" in block:
            parts.append(
                " ".join(
                    "" if value is None else str(value)
                    for value in block["cells"]
                )
            )
        elif "value" in block:
            parts.append(
                json.dumps(
                    block["value"],
                    ensure_ascii=False,
                    sort_keys=True,
                    separators=(",", ":"),
                    allow_nan=False,
                )
            )
        elif "locator" in block:
            parts.append(str(block["locator"]))
    return "\n".join(parts)


def _locators(blocks: list[dict[str, Any]]) -> tuple[str, ...]:
    return tuple(str(block["locator"]) for block in blocks)


def _enforce_block_limit(count: int, limits: DocumentLimits) -> None:
    if count > limits.maximum_blocks:
        raise ValueError("document block limit exceeded")


def _unique_json_object(pairs: list[tuple[str, object]]) -> dict[str, object]:
    output: dict[str, object] = {}
    for key, value in pairs:
        if key in output:
            raise ValueError(f"duplicate JSON object key: {key}")
        output[key] = value
    return output


def _reject_json_constant(value: str) -> object:
    raise ValueError(f"non-finite JSON number is forbidden: {value}")


def _json_pointer_token(value: str) -> str:
    return value.replace("~", "~0").replace("/", "~1")


def _json_scalar(value: object) -> str | int | float | bool | None:
    if value is None or isinstance(value, (str, int, bool)):
        return value
    if isinstance(value, float):
        if not math.isfinite(value):
            raise ValueError("non-finite spreadsheet values are forbidden")
        return value
    if isinstance(value, (date, datetime)):
        return value.isoformat()
    return str(value)


def _validate_member_name(name: str) -> None:
    path = PurePosixPath(name)
    if path.is_absolute() or ".." in path.parts or "\x00" in name:
        raise ValueError("archive path traversal detected")


def _looks_like_csv(body: bytes) -> bool:
    try:
        sample = body[:8_192].decode("utf-8-sig")
    except UnicodeDecodeError:
        return False
    if "\n" not in sample or not any(delimiter in sample for delimiter in ",;\t"):
        return False
    try:
        csv.Sniffer().sniff(sample, delimiters=",;\t")
    except csv.Error:
        return False
    return True


def _looks_like_text(body: bytes) -> bool:
    if not body:
        return True
    if b"\x00" in body[:8_192]:
        return False
    try:
        sample = body[:8_192].decode("utf-8-sig")
    except UnicodeDecodeError:
        return False
    printable = sum(character.isprintable() or character in "\r\n\t" for character in sample)
    return printable / max(1, len(sample)) >= 0.9
